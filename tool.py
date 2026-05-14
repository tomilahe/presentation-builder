"""
title: Presentation generating tool
author: Tomi Lahe
description: Loads PowerPoint templates and generates presentations from structured specs.
version: 1.0.0
licence: MIT
"""

import json
import os
import uuid
from typing import Any, Callable, Dict, List, Optional, Tuple

from pydantic import BaseModel, Field
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import ChartData
from pptx.enum.text import MSO_AUTO_SIZE

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

_CHART_TYPE_MAP: Dict[str, Any] = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
}

_PH_TITLE = 1
_PH_BODY = 2
_PH_CENTER_TITLE = 3
_PH_SUBTITLE = 4
_PH_OBJECT = 7
_PH_DATE = 16
_PH_FOOTER = 15
_PH_SLIDE_NUM = 12

_META_PH_TYPES = {_PH_DATE, _PH_FOOTER, _PH_SLIDE_NUM}

_PH_TYPE_LABELS: Dict[int, str] = {
    _PH_TITLE: "TITLE",
    _PH_BODY: "BODY",
    _PH_CENTER_TITLE: "CENTER_TITLE",
    _PH_SUBTITLE: "SUBTITLE",
    _PH_OBJECT: "OBJECT",
}


# ---------------------------------------------------------------------------
# Pydantic models
# ---------------------------------------------------------------------------


class ChartSpec(BaseModel):
    chart_type: str = "column"
    categories: List[str]
    series: Dict[str, List[float]]
    left: float = 0.7
    top: float = 1.6
    width: float = 11.9
    height: float = 5.5
    title: Optional[str] = None
    has_legend: bool = True
    legend_position: Optional[str] = "bottom"


class TableSpec(BaseModel):
    data: List[List[str]]
    has_header: bool = True
    left: float = 0.7
    top: float = 1.6
    width: float = 11.9
    col_widths: Optional[List[float]] = None


class SlideSpec(BaseModel):
    layout: int
    placeholders: Dict[str, str] = Field(default_factory=dict)
    chart: Optional[ChartSpec] = None
    table: Optional[TableSpec] = None


class PresentationSpec(BaseModel):
    slides: List[SlideSpec] = Field(min_length=1)


# ---------------------------------------------------------------------------
# Tool Class
# ---------------------------------------------------------------------------


class Tools:

    class Valves(BaseModel):
        # Output directory, public base URL for downloads, and template search path.
        STATIC_DIR: str = "/app/backend/open_webui/static/presentations"
        STATIC_URL: str = "https://thesis-test.duckdns.org/static/presentations"
        TEMPLATES_DIR: str = "/app/backend/data/templates"

    class UserValves(BaseModel):
        # Request-scoped settings (e.g. which template file to open).
        template: str = Field(
            default="icosagen_standard",
            description="Presentation template to use.",
            json_schema_extra={
                "input": {
                    "type": "select",
                    "options": "get_template_options",
                }
            },
        )

        @classmethod
        def get_template_options(cls, __user__=None) -> list[dict]:
            """List .pptx templates under Valves.TEMPLATES_DIR as value/label pairs."""
            templates_dir = Tools.Valves().TEMPLATES_DIR
            if not os.path.isdir(templates_dir):
                return [{"value": "icosagen_standard", "label": "icosagen_standard"}]
            options = [
                {
                    "value": os.path.splitext(f)[0],
                    "label": os.path.splitext(f)[0].replace("_", " "),
                }
                for f in sorted(os.listdir(templates_dir))
                if f.endswith(".pptx")
            ]
            return (
                options
                if options
                else [{"value": "icosagen_standard", "label": "icosagen_standard"}]
            )

    def __init__(self):
        self.valves = self.Valves()
        self.user_valves = self.UserValves()
        os.makedirs(self.valves.STATIC_DIR, exist_ok=True)

    # -----------------------------------------------------------------------
    # Template path
    # -----------------------------------------------------------------------

    def _get_template_path(self, __user__: Optional[dict]) -> str:
        """Absolute path to the template .pptx (Valves.TEMPLATES_DIR + user template)."""
        user_valves = (__user__ or {}).get("valves")
        template_name = getattr(user_valves, "template", "icosagen_standard")
        return os.path.join(self.valves.TEMPLATES_DIR, f"{template_name}.pptx")

    # -----------------------------------------------------------------------
    # Template manifest
    # -----------------------------------------------------------------------

    async def get_template_manifest(
        self,
        __user__: Optional[dict] = None,
    ) -> Dict[str, Any]:
        """
        Inspect the active template and return slide dimensions, layout indices,
        placeholder indices and types, and supported chart types for building a
        PresentationSpec.
        """
        template_path = self._get_template_path(__user__)

        if not os.path.exists(template_path):
            available = (
                [
                    os.path.splitext(f)[0]
                    for f in os.listdir(self.valves.TEMPLATES_DIR)
                    if f.endswith(".pptx")
                ]
                if os.path.isdir(self.valves.TEMPLATES_DIR)
                else []
            )
            return {
                "ok": False,
                "error": (
                    f"Template '{os.path.basename(template_path)}' not found. "
                    f"Available: {available}"
                ),
            }

        prs = PPTXPresentation(template_path)
        layouts = []

        for i, layout in enumerate(prs.slide_layouts):
            placeholders = []
            body_count = 0

            for shape in layout.shapes:
                if not getattr(shape, "is_placeholder", False):
                    continue
                pf = shape.placeholder_format
                pht = int(pf.type)
                if pht in _META_PH_TYPES:
                    continue

                if pht in (_PH_TITLE, _PH_CENTER_TITLE):
                    key = "title"
                elif pht == _PH_SUBTITLE:
                    key = "subtitle"
                elif pht in (_PH_BODY, _PH_OBJECT):
                    key = "body" if body_count == 0 else f"body{body_count + 1}"
                    body_count += 1
                else:
                    key = str(pf.idx)

                placeholders.append(
                    {
                        "ph_idx": pf.idx,
                        "ph_type": _PH_TYPE_LABELS.get(pht, str(pht)),
                        "key": key,
                        "left_in": round(shape.left.inches, 2),
                        "top_in": round(shape.top.inches, 2),
                        "width_in": round(shape.width.inches, 2),
                        "height_in": round(shape.height.inches, 2),
                    }
                )

            layouts.append(
                {
                    "index": i,
                    "name": layout.name,
                    "placeholders": placeholders,
                }
            )

        return {
            "ok": True,
            "slide_width_inches": prs.slide_width.inches,
            "slide_height_inches": prs.slide_height.inches,
            "num_layouts": len(layouts),
            "layouts": layouts,
            "supported_chart_types": sorted(_CHART_TYPE_MAP.keys()),
        }

    # -----------------------------------------------------------------------
    # Render presentation
    # -----------------------------------------------------------------------

    async def render_presentation(
        self,
        spec: Any,
        filename_prefix: str = "presentation",
        __user__: Optional[dict] = None,
        __event_emitter__: Optional[Callable] = None,
    ) -> Dict[str, Any]:
        """
        Render a PresentationSpec to a .pptx file and return a download URL.

        Only call this after the user has explicitly confirmed the outline
        produced from get_template_manifest. Never call this in the same
        turn as get_template_manifest.

        Each slide uses SlideSpec.layout (manifest layout index), placeholders
        (keys resolved to placeholder indices), and optionally chart or table.
        If both chart and table are present on a slide, only the chart is drawn.
        """

        async def emit_status(
            msg: str, done: bool = False, hidden: bool = False
        ) -> None:
            if __event_emitter__:
                await __event_emitter__(
                    {
                        "type": "status",
                        "data": {"description": msg, "done": done, "hidden": hidden},
                    }
                )

        async def emit_notification(msg: str, level: str = "info") -> None:
            if __event_emitter__:
                await __event_emitter__(
                    {"type": "notification", "data": {"type": level, "content": msg}}
                )

        # Validation
        await emit_status("Validating spec...")
        try:
            if isinstance(spec, str):
                spec = json.loads(spec)
            parsed = PresentationSpec.model_validate(spec)
        except Exception as exc:
            return {
                "status": "error",
                "message": f"Schema error: {exc}",
                "file_url": None,
            }

        # Template
        await emit_status("Loading template...")
        template_path = self._get_template_path(__user__)
        if not os.path.exists(template_path):
            await emit_notification(f"Template file missing: {template_path}", "error")
            return {
                "status": "error",
                "message": f"Template file missing: {template_path}",
                "file_url": None,
            }

        prs = PPTXPresentation(template_path)

        # Layout introspection
        max_li = len(prs.slide_layouts) - 1
        ph_maps = self._get_layout_ph_maps(prs)

        # Normalisation
        norm_slides, errors, warnings = self._normalize_spec(parsed, ph_maps, max_li)

        if errors:
            await emit_status(f"Spec rejected — {len(errors)} error(s)", done=True)
            return {
                "status": "error",
                "message": "; ".join(errors),
                "warnings": warnings,
                "file_url": None,
            }

        # Rendering
        total = len(norm_slides)
        await emit_status(f"Rendering {total} slide(s)...")

        for i, sd in enumerate(norm_slides, start=1):
            layout_idx = sd["layout"]

            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

            ph_shape_map: Dict[int, Any] = {
                shape.placeholder_format.idx: shape
                for shape in slide.shapes
                if getattr(shape, "is_placeholder", False)
            }

            if sd.get("placeholders"):
                self._fill_placeholders(ph_shape_map, sd["placeholders"])
            populated = set(sd.get("placeholders", {}).keys())
            for idx, shape in ph_shape_map.items():
                if str(idx) not in populated and hasattr(shape, "text_frame"):
                    shape.text_frame.clear()
            if sd.get("chart"):
                self._add_chart(slide, ChartSpec(**sd["chart"]))

            if sd.get("table"):
                self._add_table(slide, TableSpec(**sd["table"]))

        await emit_status("Saving presentation...")
        filename = f"{filename_prefix}_{uuid.uuid4().hex[:8]}.pptx"
        out_path = os.path.join(self.valves.STATIC_DIR, filename)
        prs.save(out_path)

        url = self.valves.STATIC_URL.rstrip("/") + "/" + filename

        summary_parts = [f"{total} slide(s)"]
        if warnings:
            summary_parts.append(f"{len(warnings)} warning(s)")

        await emit_status(f"Done — {', '.join(summary_parts)}", done=True)
        await emit_notification("Presentation ready", "success")

        result: Dict[str, Any] = {
            "status": "ok",
            "message": f"Presentation ready: {url}",
            "file_url": url,
            "slide_count": total,
        }

        if warnings:
            result["warnings"] = warnings
        return result

    # -----------------------------------------------------------------------
    # Helper for normalization
    # -----------------------------------------------------------------------

    def _normalize_spec(
        self,
        parsed: PresentationSpec,
        ph_maps: Dict[int, Dict[int, Dict]],
        max_li: int,
    ) -> Tuple[List[Dict], List[str], List[str]]:
        """
        Shared validation + normalization used by render_presentation.
        Returns (norm_slides, errors, warnings).
        """
        errors: List[str] = []
        warnings: List[str] = []
        norm_slides: List[Dict] = []

        for idx, slide in enumerate(parsed.slides):
            layout = max(0, min(int(slide.layout), max_li))
            if layout != slide.layout:
                warnings.append(
                    f"Slide {idx}: layout {slide.layout} clamped to {layout}"
                )

            available = ph_maps.get(layout, {})
            resolved: Dict[str, str] = {}

            for key, value in slide.placeholders.items():
                ph_idx = self._resolve_key(key, available)
                if ph_idx is None:
                    warnings.append(
                        f"Slide {idx}: key '{key}' not found in layout {layout} "
                        f"(available ph_idx: {sorted(available.keys())}). Dropped."
                    )
                    continue
                if len(str(value)) > 400:
                    warnings.append(
                        f"Slide {idx}: key '{key}' has {len(str(value))} chars — "
                        "will auto-shrink. Consider splitting across slides."
                    )
                resolved[str(ph_idx)] = str(value)

            # chart validation
            norm_chart = None
            if slide.chart:
                c = slide.chart
                if c.chart_type not in _CHART_TYPE_MAP:
                    errors.append(
                        f"Slide {idx}: unknown chart_type '{c.chart_type}'. "
                        f"Supported: {sorted(_CHART_TYPE_MAP.keys())}"
                    )
                elif not c.categories:
                    errors.append(f"Slide {idx}: chart has no categories")
                elif not c.series:
                    errors.append(f"Slide {idx}: chart has no series")
                else:
                    length_ok = all(
                        len(v) == len(c.categories) for v in c.series.values()
                    )
                    if not length_ok:
                        for sname, vals in c.series.items():
                            if len(vals) != len(c.categories):
                                errors.append(
                                    f"Slide {idx}: series '{sname}' length "
                                    f"{len(vals)} != categories {len(c.categories)}"
                                )
                    else:
                        if c.top < 1.4:
                            warnings.append(
                                f'Slide {idx}: chart top={c.top}" may overlap title '
                                f'(recommend >=1.6")'
                            )
                        norm_chart = c.model_dump()

            # table validation
            norm_table = None
            if slide.table:
                t = slide.table
                if not t.data:
                    errors.append(f"Slide {idx}: table has no data")
                elif len(t.data) < 2:
                    warnings.append(
                        f"Slide {idx}: table has only {len(t.data)} row(s) — "
                        "consider using a placeholder instead"
                    )
                else:
                    col_count = len(t.data[0])
                    bad_rows = [
                        r for r, row in enumerate(t.data) if len(row) != col_count
                    ]
                    if bad_rows:
                        errors.append(
                            f"Slide {idx}: table rows {bad_rows} have different "
                            f"column counts (expected {col_count})"
                        )
                    elif t.col_widths is not None:
                        if len(t.col_widths) != col_count:
                            errors.append(
                                f"Slide {idx}: col_widths has {len(t.col_widths)} "
                                f"entries but table has {col_count} columns"
                            )
                        elif sum(t.col_widths) > t.width + 0.01:
                            errors.append(
                                f"Slide {idx}: col_widths sum "
                                f'({sum(t.col_widths):.2f}") exceeds table '
                                f'width ({t.width}")'
                            )
                        else:
                            norm_table = t.model_dump()
                    else:
                        norm_table = t.model_dump()

                if slide.chart and norm_table:
                    warnings.append(
                        f"Slide {idx}: has both a chart and a table — "
                        "only the chart will be rendered"
                    )
                    norm_table = None

            norm_slides.append(
                {
                    "layout": layout,
                    "placeholders": resolved,
                    "chart": norm_chart,
                    "table": norm_table,
                }
            )

        return norm_slides, errors, warnings

    # -----------------------------------------------------------------------
    # Helper functions for layout introspection
    # -----------------------------------------------------------------------

    def _get_layout_ph_maps(self, prs) -> Dict[int, Dict[int, Dict]]:
        """
        Build {layout_idx: {ph_idx: {ph_type, …}}} by reading layout shapes
        directly — no slide instances are created.
        """
        maps: Dict[int, Dict[int, Dict]] = {}
        for li, layout in enumerate(prs.slide_layouts):
            ph_map: Dict[int, Dict] = {}
            for shape in layout.shapes:
                if not getattr(shape, "is_placeholder", False):
                    continue
                pf = shape.placeholder_format
                pht = int(pf.type)
                if pht in _META_PH_TYPES:
                    continue
                ph_map[pf.idx] = {
                    "ph_type": pht,
                    "left_in": round(shape.left.inches, 2),
                    "top_in": round(shape.top.inches, 2),
                    "width_in": round(shape.width.inches, 2),
                    "height_in": round(shape.height.inches, 2),
                }
            maps[li] = ph_map
        return maps

    def _resolve_key(self, key: str, available: Dict[int, Dict]) -> Optional[int]:
        """Resolve a placeholder key to a ph_idx integer."""
        k = key.strip()

        if k.isdigit():
            idx = int(k)
            return idx if idx in available else None

        kl = k.lower()

        def first_of(*types: int) -> Optional[int]:
            for ph_type in types:
                for idx in sorted(available):
                    if available[idx]["ph_type"] == ph_type:
                        return idx
            return None

        def nth_body(n: int) -> Optional[int]:
            count = 0
            for idx in sorted(available):
                if available[idx]["ph_type"] in (_PH_BODY, _PH_OBJECT):
                    if count == n:
                        return idx
                    count += 1
            return None

        if kl == "title":
            return first_of(_PH_TITLE, _PH_CENTER_TITLE)
        if kl == "subtitle":
            return first_of(_PH_SUBTITLE)
        if kl in ("body", "body1", "content", "content1"):
            return nth_body(0)
        if kl in ("body2", "content2"):
            return nth_body(1)
        if kl in ("body3", "content3"):
            return nth_body(2)

        return None

    # -----------------------------------------------------------------------
    # Private helpers — text rendering
    # -----------------------------------------------------------------------

    def _fill_placeholders(
        self, ph_shape_map: Dict[int, Any], placeholders: Dict[str, str]
    ) -> None:
        """Fill numeric-index placeholders from text.
        Newlines become paragraphs.
        """
        for key, value in placeholders.items():
            if not key.isdigit():
                continue
            shape = ph_shape_map.get(int(key))
            if shape is None or not hasattr(shape, "text_frame"):
                continue
            tf = shape.text_frame
            tf.clear()
            lines = value.split("\n") if value else [""]
            tf.paragraphs[0].add_run().text = lines[0]
            for line in lines[1:]:
                tf.add_paragraph().add_run().text = line
            shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # -----------------------------------------------------------------------
    # Private helpers — image, chart, table
    # -----------------------------------------------------------------------

    def _add_chart(self, slide, spec: ChartSpec) -> None:
        """Insert a chart from ChartSpec.
        Skip if chart_type is unsupported.
        """
        chart_type = _CHART_TYPE_MAP.get(spec.chart_type)
        if not chart_type:
            return

        cd = ChartData()
        cd.categories = spec.categories
        for name, values in spec.series.items():
            cd.add_series(str(name), values)

        shape = slide.shapes.add_chart(
            chart_type,
            Inches(spec.left),
            Inches(spec.top),
            Inches(spec.width),
            Inches(spec.height),
            cd,
        )
        chart = shape.chart

        if spec.title:
            chart.has_title = True
            chart.chart_title.text_frame.text = spec.title
        else:
            chart.has_title = False

        chart.has_legend = spec.has_legend
        if spec.has_legend and spec.legend_position:
            _pos_map = {
                "right": XL_LEGEND_POSITION.RIGHT,
                "left": XL_LEGEND_POSITION.LEFT,
                "top": XL_LEGEND_POSITION.TOP,
                "bottom": XL_LEGEND_POSITION.BOTTOM,
                "corner": XL_LEGEND_POSITION.CORNER,
            }
            pos = _pos_map.get(spec.legend_position.lower())
            if pos is not None:
                chart.legend.position = pos
                chart.legend.include_in_layout = False

    def _add_table(self, slide, spec: TableSpec) -> None:
        """
        Insert a native PowerPoint table from TableSpec.
        When has_header is True, the first row uses bold text; other rows are
        plain. Column widths default to an even split of spec.width, or use
        col_widths when set. Cell fonts follow the template table style (no
        explicit font is applied here).
        """
        if not spec.data:
            return

        rows = len(spec.data)
        cols = len(spec.data[0])

        # Distribute column widths evenly if not specified.
        if spec.col_widths:
            col_widths_emu = [Inches(w) for w in spec.col_widths]
        else:
            col_w = spec.width / cols
            col_widths_emu = [Inches(col_w)] * cols

        # python-pptx requires a total width and height for add_table;
        # height is estimated at 0.4 inches per row as a sensible default
        # PowerPoint will expand rows to fit content automatically.
        row_height = Inches(0.4)
        total_height = row_height * rows

        graphic_frame = slide.shapes.add_table(
            rows,
            cols,
            Inches(spec.left),
            Inches(spec.top),
            Inches(spec.width),
            total_height,
        )
        table = graphic_frame.table

        # Apply explicit column widths.
        for ci, width_emu in enumerate(col_widths_emu):
            table.columns[ci].width = width_emu

        # Fill cell text. Header row gets bold runs, body rows are plain.
        for ri, row_data in enumerate(spec.data):
            is_header = spec.has_header and ri == 0
            for ci, cell_text in enumerate(row_data):
                cell = table.cell(ri, ci)
                tf = cell.text_frame
                tf.clear()
                para = tf.paragraphs[0]
                run = para.add_run()
                run.text = str(cell_text)
                if is_header:
                    run.font.bold = True

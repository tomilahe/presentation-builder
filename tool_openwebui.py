from __future__ import annotations

from pydantic import BaseModel, Field
from typing import Any, Dict, Optional, List
from datetime import datetime
from pptx import Presentation as PPTXPresentation
import os
import json


class GenerationResult(BaseModel):
    status: str
    message: str
    file_url: str
    mimetype: str
    slides: Optional[Any] = None


class Tools:
    class Valves(BaseModel):
        STATIC_DIR: str = Field("/app/static")
        STATIC_URL: str = Field("https://icosagenai.hpc.ut.ee/static/")
        template_path: str = Field("/app/backend/data/icosagen-template.pptx")

    def __init__(self):
        self.valves = self.Valves()
        os.makedirs(self.valves.STATIC_DIR, exist_ok=True)

    # -------------------------
    # JSON -> PPTX
    # -------------------------

    def _build_and_save_from_json(
        self,
        data: Dict[str, Any],
        filename_prefix: str,
        template_path: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        data schema:

        {
          "slides": [
            {
              "layout": <int>,
              "placeholders": {
                "<Placeholder Name>": "<text or \\n-separated bullets>"
              }
            }
          ]
        }
        """
        final_template = template_path or self.valves.template_path
        if not os.path.exists(final_template):
            raise FileNotFoundError(f"Template not found: {final_template}")

        prs = PPTXPresentation(final_template)
        slides_data = data.get("slides", [])

        if not isinstance(slides_data, list) or not slides_data:
            raise ValueError("Input must contain a non-empty 'slides' list")

        # Fill first existing slide in template
        first_slide_data = slides_data[0]
        first_placeholders = first_slide_data.get("placeholders", {})
        if prs.slides:
            first_slide = prs.slides[0]
        else:
            layout_idx = first_slide_data.get("layout", 0)
            if not isinstance(layout_idx, int):
                layout_idx = 0
            num_layouts = len(prs.slide_layouts)
            if layout_idx < 0 or layout_idx >= num_layouts:
                layout_idx = 0
            layout = prs.slide_layouts[layout_idx]
            first_slide = prs.slides.add_slide(layout)
        self._fill_slide_placeholders(first_slide, first_placeholders)

        # Add new slides for the rest
        for slide_data in slides_data[1:]:
            layout_idx = slide_data.get("layout", 1)

            # Defensive: clamp layout index
            if not isinstance(layout_idx, int):
                layout_idx = 1
            num_layouts = len(prs.slide_layouts)
            if layout_idx < 0 or layout_idx >= num_layouts:
                layout_idx = 1

            layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(layout)
            self._fill_slide_placeholders(slide, slide_data.get("placeholders", {}))

        # Save file
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"{filename_prefix}_{timestamp}.pptx"
        output_path = os.path.join(self.valves.STATIC_DIR, filename)
        prs.save(output_path)

        file_url = self.valves.STATIC_URL.rstrip("/") + "/" + filename

        return {
            "status": "success",
            "message": f"Generated: {file_url}",
            "file_url": file_url,
            "mimetype": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        }

    def _fill_slide_placeholders(self, slide, placeholders: Dict[str, Any]) -> None:
        for shape in slide.shapes:
            if not getattr(shape, "is_placeholder", False):
                continue
            name = getattr(shape, "name", None)
            if not name:
                continue
            if name in placeholders and hasattr(shape, "text_frame"):
                shape.text_frame.clear()
                shape.text_frame.text = str(placeholders[name])

    # -------------------------
    # OpenWebUI entrypoint
    # -------------------------

    async def create_presentation(
        self,
        slides: Dict[str, Any],
        __event_emitter__=None,
    ) -> GenerationResult:
        """
        slides: JSON-compatible dict with structure:

        {
          "slides": [
            {
              "layout": <int>,
              "placeholders": {
                "<Placeholder Name>": "<text or \\n-separated bullets>"
              }
            }
          ]
        }
        """
        try:
            result = self._build_and_save_from_json(
                data=slides,
                filename_prefix="presentation",
                template_path=None,
            )
            return GenerationResult(
                status="ok",
                message=result["message"],
                file_url=result["file_url"],
                mimetype=result["mimetype"],
                slides=slides,
            )
        except Exception as e:
            return GenerationResult(
                status="error",
                message=str(e),
                file_url="",
                mimetype=(
                    "application/"
                    "vnd.openxmlformats-officedocument.presentationml.presentation"
                ),
                slides=None,
            )
